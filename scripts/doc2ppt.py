#!/usr/bin/env python3
"""
via54Design — Document / Image → PPT Pipeline
Accepts: .docx, .md, .txt, .pptx, .png/.jpg (image)
Outputs: structured PPT framework + content
"""
import json, sys, os, re, textwrap
from pathlib import Path

try:
    from docx import Document
except: Document = None
try:
    from pptx import Presentation
except: Presentation = None

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from img2prompt import analyze_image, build_prompt_from_analysis

def extract_text_from_docx(path: str) -> dict:
    """Extract text structure from .docx"""
    doc = Document(path)
    title = ""
    paragraphs = []
    headings = []
    current_section = "intro"
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading"):
            level = para.style.name.replace("Heading ", "")
            try: level = int(level)
            except: level = 2
            headings.append({"level": level, "text": text})
            current_section = text
        else:
            paragraphs.append({"section": current_section, "text": text})
    
    # Extract images
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                images.append({"name": rel.target_ref.split("/")[-1], "type": "embedded"})
            except: pass
    
    if not title and headings:
        title = headings[0]["text"]
    
    return {
        "type": "docx",
        "title": title or Path(path).stem,
        "headings": headings,
        "paragraphs": paragraphs,
        "total_paragraphs": len(paragraphs),
        "total_headings": len(headings),
        "images": images,
        "content_preview": "\n".join(p["text"][:120] for p in paragraphs[:10]),
    }

def extract_text_from_md(path: str) -> dict:
    """Extract structure from Markdown"""
    with open(path, encoding="utf-8") as f:
        content = f.read()
    
    headings = []
    paragraphs = []
    for line in content.split("\n"):
        stripped = line.strip()
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            text = stripped.lstrip("#").strip()
            headings.append({"level": level, "text": text})
        elif stripped and not stripped.startswith(("```", "---", ">", "-", "*", "|")):
            paragraphs.append(stripped[:200])
    
    # Find title (first H1)
    title = ""
    for h in headings:
        if h["level"] == 1:
            title = h["text"]
            break
    
    return {
        "type": "markdown",
        "title": title or Path(path).stem,
        "headings": headings,
        "total_headings": len(headings),
        "content_preview": "\n".join(paragraphs[:8]),
        "raw_length": len(content),
    }

def extract_from_pptx(path: str) -> dict:
    """Extract content from existing PPTX"""
    if Presentation is None:
        return {"type": "pptx", "error": "python-pptx not installed"}
    
    prs = Presentation(path)
    slides = []
    all_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_data = {"index": i + 1, "texts": [], "image_count": 0}
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_data["texts"].append(t)
                        all_text.append(t)
            if shape.shape_type == 13:  # Picture
                slide_data["image_count"] += 1
        slides.append(slide_data)
    
    # Detect structure
    has_titles = any(len(s["texts"]) > 0 for s in slides[:3])
    
    return {
        "type": "pptx",
        "title": Path(path).stem,
        "total_slides": len(slides),
        "total_images": sum(s["image_count"] for s in slides),
        "total_text_blocks": sum(len(s["texts"]) for s in slides),
        "slides": slides[:20],  # first 20 slides
        "content_preview": "\n".join(all_text[:15]),
    }

def extract_from_image(path: str) -> dict:
    """Extract visual analysis from image for PPT context"""
    analysis = analyze_image(path)
    prompt = build_prompt_from_analysis(analysis)
    
    suggested_topics = []
    moods = analysis.get("suggested_moods", [])
    if "vibrant" in analysis.get("colorfulness", ""):
        suggested_topics = ["品牌展示", "产品亮点", "创意概念"]
    elif "muted" in analysis.get("colorfulness", ""):
        suggested_topics = ["专业报告", "数据分析", "行业洞察"]
    else:
        suggested_topics = ["综合提案", "方案介绍", "案例分享"]
    
    return {
        "type": "image",
        "title": Path(path).stem,
        "analysis": analysis,
        "generated_prompt": prompt,
        "suggested_topics": suggested_topics,
        "visual_context": {
            "colors": [c["hex"] for c in (analysis.get("dominant_colors") or [])[:5]],
            "brightness": analysis.get("brightness_label"),
            "moods": moods,
            "styles": analysis.get("suggested_styles", []),
        }
    }

def extract_content(path: str) -> dict:
    """Auto-detect file type and extract content"""
    ext = Path(path).suffix.lower()
    
    if ext in (".docx",):
        if Document is None:
            return {"error": "python-docx not installed. Run: pip install python-docx"}
        return extract_text_from_docx(path)
    elif ext in (".md", ".markdown"):
        return extract_text_from_md(path)
    elif ext in (".txt",):
        with open(path, encoding="utf-8") as f:
            text = f.read()
        return {
            "type": "text",
            "title": Path(path).stem,
            "content_preview": text[:800],
            "raw_length": len(text),
            "total_lines": text.count("\n") + 1,
        }
    elif ext in (".pptx",):
        return extract_from_pptx(path)
    elif ext in (".png", ".jpg", ".jpeg", ".webp", ".gif"):
        return extract_from_image(path)
    else:
        return {"error": f"unsupported format: {ext}"}

def generate_ppt_framework(content: dict, user_prompt: str = "") -> dict:
    """Generate PPT framework from extracted content."""
    ctype = content.get("type", "unknown")
    title = content.get("title", "演示文稿")
    slides = []
    
    # Cover slide
    slides.append({
        "type": "cover",
        "title": title,
        "subtitle": user_prompt or "via54Design 生成",
        "mood": "inspiring"
    })
    
    if ctype == "docx":
        # Build from headings
        headings = content.get("headings", [])
        for h in headings[:15]:
            slides.append({
                "type": "content",
                "title": h["text"],
                "subtitle": f"H{h['level']} · 章节页",
                "mood": "informative",
                "layout": "hero-split-16-9",
            })
        # Summary
        slides.append({"type": "summary", "title": "总结", "mood": "confident"})
        
    elif ctype == "markdown":
        headings = content.get("headings", [])
        for h in headings[:12]:
            slides.append({
                "type": "content",
                "title": h["text"],
                "subtitle": f"Level {h['level']}",
                "mood": "informative",
            })
        slides.append({"type": "summary", "title": "总结与展望", "mood": "inspiring"})
        
    elif ctype == "pptx":
        existing = content.get("slides", [])
        for s in existing[:15]:
            texts = s.get("texts", [])
            slide_title = texts[0] if texts else f"第{s['index']}页"
            slides.append({
                "type": "content",
                "title": slide_title,
                "subtitle": f"Slide {s['index']}",
                "mood": "informative",
                "has_images": s.get("image_count", 0) > 0,
            })
        slides.append({"type": "summary", "title": "总结", "mood": "confident"})
        
    elif ctype == "image":
        vc = content.get("visual_context", {})
        moods = vc.get("moods", ["professional"])
        topics = content.get("suggested_topics", ["方案介绍"])
        
        slides.append({
            "type": "section",
            "title": "视觉灵感",
            "subtitle": f"风格: {' · '.join(vc.get('styles', ['professional']))} | 情绪: {' · '.join(moods)}",
            "mood": moods[0] if moods else "inspiring",
        })
        for topic in topics[:4]:
            slides.append({
                "type": "content",
                "title": topic,
                "mood": "informative",
                "image_hint": content.get("generated_prompt", "")[:100],
            })
        slides.append({"type": "summary", "title": "下一步", "mood": "hopeful"})
        
    else:  # text / unknown
        preview = content.get("content_preview", "")
        sentences = [s.strip() for s in re.split(r'[。！？\n]', preview) if len(s.strip()) > 10]
        for s in sentences[:8]:
            slides.append({
                "type": "content",
                "title": s[:60],
                "mood": "informative",
            })
        slides.append({"type": "summary", "title": "总结", "mood": "confident"})
    
    return {
        "title": title,
        "type": ctype,
        "total_slides": len(slides),
        "slides": slides,
        "content_info": {
            "paragraphs": content.get("total_paragraphs", 0),
            "headings": content.get("total_headings", 0),
            "images": content.get("total_images", content.get("images", 0)),
        },
        "recommended_command": (
            f"via54 generate --layout hero-split-16-9 "
            f"--color ink-wash --font ming-hei-editorial "
            f"--title \"{title[:40]}\" --presentation"
        ),
        "user_guidance": generate_guidance(ctype, content),
    }

def generate_guidance(ctype: str, content: dict) -> str:
    """Generate guidance for user on next steps."""
    guides = {
        "docx": "📄 检测到 Word 文档。已提取章节结构，建议逐章确认标题，\n"
                "然后选择配色方案和字体，点击「生成完整演示」。",
        "markdown": "📝 检测到 Markdown 文件。已解析标题层级，\n"
                    "建议补充每节的核心要点（每节 3-5 个 bullet），然后生成。",
        "pptx": "📊 检测到已有 PPTX 文件。已提取每页文本和图片信息，\n"
                "可以选择保留原始风格或应用 via54 模板重新设计。",
        "image": "🖼️ 检测到图片。已分析视觉特征（色彩/亮度/情绪），\n"
                 f"建议基于「{' · '.join(content.get('suggested_topics', ['方案']))}」方向补充文案。",
        "text": "📃 检测到纯文本。已自动分段，建议标注各段标题，\n"
                "并补充关键数据和图片需求。",
    }
    return guides.get(ctype, "文件已分析完成，请确认内容并继续。")

def story2ppt(path: str, user_prompt: str = "", mode: str = "auto") -> dict:
    """Main entry: file → content → PPT framework."""
    content = extract_content(path)
    if "error" in content:
        return content
    
    framework = generate_ppt_framework(content, user_prompt)
    framework["source_file"] = path
    framework["user_prompt"] = user_prompt
    return framework

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Document → PPT Pipeline")
    parser.add_argument("path", help="Path to file (.docx/.md/.txt/.pptx/.png/.jpg)")
    parser.add_argument("--prompt", default="", help="User description / topic")
    args = parser.parse_args()
    
    result = story2ppt(args.path, args.prompt)
    print(json.dumps(result, indent=2, ensure_ascii=False))
